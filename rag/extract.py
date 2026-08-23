"""Get text out of a document, and refuse rather than pretend when we cannot.

Format is decided by magic bytes, never by the filename or the client's
content-type: both are attacker-controlled and both are routinely wrong even
when nobody is being hostile.

Every parser here is pure Python with aarch64 wheels, so the image needs no
compiler and the arm64 build stays a copy rather than a compile.
"""
import io
import zipfile

from .config import (CHARS_PER_TOKEN, MAX_EXTRACT_CHARS, MAX_EXTRACT_PAGES,
                     MIN_CHARS_PER_PAGE)


class ExtractError(Exception):
    """Ingest failed for a reason the caller can act on."""


def sniff(data: bytes, filename: str | None = None) -> str:
    head = data[:8]
    if head.startswith(b"%PDF"):
        return "application/pdf"
    if head.startswith(b"PK\x03\x04"):
        try:
            names = set(zipfile.ZipFile(io.BytesIO(data)).namelist())
        except zipfile.BadZipFile:
            raise ExtractError("looks like a zip but could not be opened")
        if "word/document.xml" in names:
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if any(n.startswith("ppt/slides/slide") for n in names):
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if "xl/workbook.xml" in names:
            raise ExtractError("spreadsheets are not supported")
        raise ExtractError("unrecognised Office/zip container")
    sample = data[:1024].lstrip().lower()
    if sample.startswith(b"<!doctype html") or sample.startswith(b"<html"):
        return "text/html"
    if filename and filename.lower().endswith((".md", ".markdown")):
        return "text/markdown"
    if filename and filename.lower().endswith((".html", ".htm")):
        return "text/html"
    return "text/plain"


def _pdf(data: bytes) -> tuple[list[dict], dict]:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    total_pages = len(reader.pages)
    blocks, chars = [], 0
    read_pages = min(total_pages, MAX_EXTRACT_PAGES)
    for i in range(read_pages):
        try:
            text = reader.pages[i].extract_text() or ""
        except Exception:
            text = ""
        chars += len(text)
        if text.strip():
            blocks.append({"text": text, "page": i + 1, "heading": None})
        if chars > MAX_EXTRACT_CHARS:
            break
    # A scanned PDF parses perfectly and yields almost nothing. Indexing it
    # would produce a document that reports ready and answers every question
    # with silence -- so it fails here, where the reason can still be explained.
    if read_pages and chars / read_pages < MIN_CHARS_PER_PAGE:
        raise ExtractError(
            "almost no text could be extracted (%d chars over %d pages) -- this "
            "looks like a scanned or image-only PDF, which needs OCR before it "
            "can be indexed" % (chars, read_pages))
    return blocks, {"total_pages": total_pages}


def _docx(data: bytes) -> tuple[list[dict], dict]:
    import docx
    doc = docx.Document(io.BytesIO(data))
    blocks, heading = [], None
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        if (para.style is not None and para.style.name
                and para.style.name.lower().startswith("heading")):
            heading = text
        blocks.append({"text": text, "page": None, "heading": heading})
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                blocks.append({"text": " | ".join(cells), "page": None,
                               "heading": heading})
    # DOCX has no page count until something renders it -- python-docx cannot
    # tell us, and guessing one would make the "20 pages" limit a lie for this
    # format. None is honest; the token budget still applies.
    return blocks, {"total_pages": None}


def _pptx(data: bytes) -> tuple[list[dict], dict]:
    from pptx import Presentation
    pres = Presentation(io.BytesIO(data))
    blocks = []
    slides = list(pres.slides)
    for i, slide in enumerate(slides[:MAX_EXTRACT_PAGES]):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if parts:
            # One block per slide: a slide is already the natural chunk, so
            # nothing is gained by cutting across them.
            blocks.append({"text": "\n\n".join(parts), "page": i + 1,
                           "heading": parts[0][:120] if parts else None})
    return blocks, {"total_pages": len(slides)}


def _html(data: bytes) -> tuple[list[dict], dict]:
    from selectolax.parser import HTMLParser
    tree = HTMLParser(data.decode("utf-8", "replace"))
    for tag in ("script", "style", "noscript", "svg"):
        for node in tree.css(tag):
            node.decompose()
    body = tree.body or tree.root
    text = body.text(separator="\n") if body else ""
    title = tree.css_first("title")
    return ([{"text": text, "page": None, "heading": None}] if text.strip() else [],
            {"total_pages": None,
             "title": title.text().strip() if title else None})


def _plain(data: bytes) -> tuple[list[dict], dict]:
    text = data.decode("utf-8", "replace")
    blocks, heading = [], None
    for para in text.split("\n\n"):
        stripped = para.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            heading = stripped.lstrip("#").strip() or heading
        blocks.append({"text": stripped, "page": None, "heading": heading})
    return blocks, {"total_pages": None}


_PARSERS = {
    "application/pdf": _pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _pptx,
    "text/html": _html,
    "text/markdown": _plain,
    "text/plain": _plain,
}


def extract(data: bytes, filename: str | None = None) -> dict:
    """bytes -> {mime, blocks, text, total_pages, total_tokens, title}."""
    mime = sniff(data, filename)
    parser = _PARSERS.get(mime)
    if parser is None:
        raise ExtractError("unsupported format: %s" % mime)
    blocks, meta = parser(data)
    if not blocks:
        raise ExtractError("no text could be extracted from this document")
    text = "\n\n".join(b["text"] for b in blocks)
    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS]
    return {
        "mime": mime,
        "blocks": blocks,
        "text": text,
        "total_pages": meta.get("total_pages"),
        "total_tokens": max(1, len(text) // CHARS_PER_TOKEN),
        "title": meta.get("title") or (filename or None),
    }
